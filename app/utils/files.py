"""
app/utils/files.py
==================
File-handling utilities:
  - Fenced code-block extraction
  - ZIP creation from code blocks
  - ZIP/RAR archive reading
  - Universal file-content extraction (PDF, DOCX, XLSX, CSV, images, ...)
  - Prompt formatter for multi-file uploads

Image files (.png .jpg .jpeg .gif .bmp .webp .tiff .tif) are no longer
processed with pytesseract OCR. Instead, they are base64-encoded and
returned as a vision payload dict:

  {"__vision__": True, "b64": <base64>, "media_type": <mime>, "meta": <str>}

The caller (app/routes/chat.py) detects this sentinel and routes the
request to ask_groq_vision() instead of the text pipeline.
"""

import base64
import csv
import io
import os
import re
import time
import zipfile

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

MAX_UPLOAD_SIZE: int = 10 * 1024 * 1024  # 10 MB

LANG_EXTENSIONS: dict[str, str] = {
    "python": ".py", "py": ".py",
    "javascript": ".js", "js": ".js",
    "typescript": ".ts", "ts": ".ts",
    "html": ".html", "css": ".css",
    "java": ".java", "c": ".c",
    "cpp": ".cpp", "c++": ".cpp", "csharp": ".cs",
    "go": ".go", "rust": ".rs", "ruby": ".rb",
    "php": ".php", "sql": ".sql", "swift": ".swift",
    "kotlin": ".kt", "dart": ".dart", "r": ".r",
    "bash": ".sh", "sh": ".sh", "shell": ".sh",
    "json": ".json", "yaml": ".yml", "yml": ".yml",
    "xml": ".xml", "markdown": ".md", "md": ".md",
    "txt": ".txt", "text": ".txt", "": ".txt",
}

CODE_EXTENSIONS: frozenset[str] = frozenset({
    ".py", ".js", ".ts", ".tsx", ".jsx", ".html", ".css", ".java",
    ".c", ".cpp", ".h", ".cs", ".go", ".rs", ".rb", ".php", ".sql",
    ".swift", ".kt", ".dart", ".r", ".sh", ".json", ".yaml", ".yml",
    ".xml", ".md", ".txt", ".toml", ".cfg", ".ini", ".env", ".gitignore",
    ".dockerfile", ".makefile",
})

IMAGE_EXTENSIONS: frozenset[str] = frozenset({
    ".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".tiff", ".tif",
})

_IMAGE_MIME: dict[str, str] = {
    ".png": "image/png",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".gif": "image/gif",
    ".bmp": "image/bmp",
    ".webp": "image/webp",
    ".tiff": "image/tiff",
    ".tif": "image/tiff",
}

_CODE_BLOCK_RE = re.compile(r"```(\w*)\n(.*?)```", re.DOTALL)

_FILE_COMMENT_RE = re.compile(
    r"^(?:#|//|/\s*\*\*)?\s*File:\s*(.+)$",
    re.IGNORECASE,
)

# ===========================================================================
# Code-block extraction & ZIP generation
# ===========================================================================

def extract_code_blocks(text: str) -> list[dict]:
    """
    Extract all fenced code blocks from Markdown text.

    Returns a list of dicts with keys language, code, filename.
    The filename is parsed from an optional # File: path comment on the
    first line of the block.
    """
    blocks = []
    for lang, content in _CODE_BLOCK_RE.findall(text):
        content = content.strip()
        filename = None
        lines = content.split("\n")
        if lines:
            match = _FILE_COMMENT_RE.match(lines[0].strip())
            if match:
                filename = match.group(1).strip()
                content = "\n".join(lines[1:]).strip()
        blocks.append({
            "language": (lang or "txt").lower(),
            "code": content,
            "filename": filename,
        })
    return blocks

def save_code_as_zip(code_blocks: list[dict]) -> dict | None:
    """
    Package named / substantial code blocks into a ZIP archive in /tmp.

    Returns a dict with filename, url, and files keys,
    or None if there is nothing worth zipping.
    """
    if not code_blocks:
        return None

    named_blocks: list[dict] = []
    for block in code_blocks:
        lang = block["language"]
        if lang in ("output", "terminal", "console"):
            continue
        if block.get("filename"):
            named_blocks.append(block)
            continue
        if lang in ("bash", "sh", "shell") and block["code"].count("\n") < 3:
            continue
        if block["code"].count("\n") >= 5:
            ext = LANG_EXTENSIONS.get(lang, ".txt")
            if len(code_blocks) == 1:
                block["filename"] = f"main{ext}"
            named_blocks.append(block)

    if not named_blocks:
        return None

    timestamp = int(time.time())
    zip_filename = f"zenith_code_{timestamp}.zip"
    zip_filepath = f"/tmp/{zip_filename}"
    used_names: set[str] = set()

    with zipfile.ZipFile(zip_filepath, "w", zipfile.ZIP_DEFLATED) as zf:
        for i, block in enumerate(named_blocks):
            lang = block["language"]
            ext = LANG_EXTENSIONS.get(lang, ".txt")
            fname = block.get("filename") or f"file_{i + 1}{ext}"
            if fname in used_names:
                base, fext = os.path.splitext(fname)
                fname = f"{base}_{i + 1}{fext}"
            used_names.add(fname)
            zf.writestr(fname, block["code"])

    if not used_names:
        return None

    return {
        "filename": zip_filename,
        "url": f"/download-zip/{zip_filename}",
        "files": sorted(used_names),
    }

# ===========================================================================
# Archive reading (ZIP / RAR)
# ===========================================================================

def read_archive(file_storage) -> tuple[dict | None, str | None]:
    """
    Extract code files from an uploaded .zip or .rar archive.

    Returns (files_dict, None) on success or (None, error_msg) on failure.
    """
    filename = file_storage.filename.lower()
    file_bytes = file_storage.read()

    if len(file_bytes) > MAX_UPLOAD_SIZE:
        return None, "File too large (max 10 MB)"

    _skip_names = {"dockerfile", "makefile", ".gitignore", ".env.example"}
    files: dict[str, str] = {}

    if filename.endswith(".zip"):
        try:
            with zipfile.ZipFile(io.BytesIO(file_bytes)) as zf:
                for name in zf.namelist():
                    if name.endswith("/") or "/__" in name or "/." in name:
                        continue
                    ext = os.path.splitext(name)[1].lower()
                    basename = os.path.basename(name)
                    if ext not in CODE_EXTENSIONS and basename.lower() not in _skip_names:
                        continue
                    try:
                        content = zf.read(name).decode("utf-8", errors="replace")
                        if len(content) > 50_000:
                            continue
                        files[name] = content
                    except Exception:
                        continue
        except zipfile.BadZipFile:
            return None, "Invalid zip file"

    elif filename.endswith(".rar"):
        try:
            import rarfile  # type: ignore[import]
            with rarfile.RarFile(io.BytesIO(file_bytes)) as rf:
                for name in rf.namelist():
                    if name.endswith("/") or "/__" in name or "/." in name:
                        continue
                    ext = os.path.splitext(name)[1].lower()
                    basename = os.path.basename(name)
                    if ext not in CODE_EXTENSIONS and basename.lower() not in _skip_names:
                        continue
                    try:
                        content = rf.read(name).decode("utf-8", errors="replace")
                        if len(content) > 50_000:
                            continue
                        files[name] = content
                    except Exception:
                        continue
        except ImportError:
            return None, "RAR support not available. Please upload a .zip file instead."
        except Exception:
            return None, "Invalid or corrupted RAR file"
    else:
        return None, "Unsupported format. Please upload .zip or .rar"

    if not files:
        return None, "No code files found in archive"

    return files, None

# ===========================================================================
# Universal file-content extractor
# ===========================================================================

def extract_file_content(file_storage) -> tuple:
    """
    Extract readable content from any uploaded file type.

    For image files returns a vision payload dict:
    ({"__vision__": True, "b64": ..., "media_type": ..., "meta": ...}, None)

    For all other types returns:
    (content_string, None) on success
    (None, error_message) on failure
    """
    filename: str = file_storage.filename or "unknown"
    ext = os.path.splitext(filename)[1].lower()
    file_bytes: bytes = file_storage.read()

    if len(file_bytes) > MAX_UPLOAD_SIZE:
        return None, "File too large (max 10 MB)"

    try:
        # PDF
        if ext == ".pdf":
            try:
                import pdfplumber  # type: ignore[import]
                text_parts = []
                with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
                    for i, page in enumerate(pdf.pages[:50]):
                        page_text = page.extract_text() or ""
                        if page_text.strip():
                            text_parts.append(f"--- Page {i + 1} ---\n{page_text}")
                if text_parts:
                    return _truncate("\n\n".join(text_parts), 30_000), None
                return None, "Could not extract text from PDF (may be image-based)"
            except ImportError:
                return None, "PDF support requires: pip install pdfplumber"

        # Word document
        elif ext in (".docx", ".doc"):
            try:
                from docx import Document  # type: ignore[import]
                doc = Document(io.BytesIO(file_bytes))
                paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
                if paragraphs:
                    return _truncate("\n\n".join(paragraphs), 30_000), None
                return None, "Word document appears to be empty"
            except ImportError:
                return None, "DOCX support requires: pip install python-docx"

        # Excel
        elif ext in (".xlsx", ".xls"):
            try:
                import openpyxl  # type: ignore[import]
                wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True)
                sheets_text = []
                for sheet_name in wb.sheetnames[:10]:
                    ws = wb[sheet_name]
                    rows = []
                    for row in ws.iter_rows(max_row=200, values_only=True):
                        row_str = " | ".join(str(c) if c is not None else "" for c in row)
                        if row_str.strip(" |"):
                            rows.append(row_str)
                    if rows:
                        sheets_text.append(f"--- Sheet: {sheet_name} ---\n" + "\n".join(rows))
                wb.close()
                if sheets_text:
                    return _truncate("\n\n".join(sheets_text), 30_000), None
                return None, "Excel file appears to be empty"
            except ImportError:
                return None, "Excel support requires: pip install openpyxl"

        # CSV / TSV
        elif ext in (".csv", ".tsv"):
            try:
                text = file_bytes.decode("utf-8", errors="replace")
                delimiter = "\t" if ext == ".tsv" else ","
                reader = csv.reader(text.splitlines(), delimiter=delimiter)
                rows = []
                for i, row in enumerate(reader):
                    if i >= 500:
                        rows.append("[... more rows truncated ...]")
                        break
                    rows.append(" | ".join(row))
                if rows:
                    return _truncate("\n".join(rows), 30_000), None
                return None, "CSV file appears to be empty"
            except Exception:
                return None, "Could not parse CSV file"

        # PowerPoint
        elif ext in (".pptx", ".ppt"):
            try:
                from pptx import Presentation  # type: ignore[import]
                prs = Presentation(io.BytesIO(file_bytes))
                slides_text = []
                for i, slide in enumerate(prs.slides):
                    texts = [
                        s.text for s in slide.shapes
                        if hasattr(s, "text") and s.text.strip()
                    ]
                    if texts:
                        slides_text.append(f"--- Slide {i + 1} ---\n" + "\n".join(texts))
                if slides_text:
                    return "\n\n".join(slides_text), None
                return None, "PowerPoint file appears to be empty"
            except ImportError:
                return None, "PPTX support requires: pip install python-pptx"

        # Images -- Groq Vision (base64 payload, replaces pytesseract OCR)
        elif ext in IMAGE_EXTENSIONS:
            try:
                from PIL import Image as PILImage  # type: ignore[import]
                img = PILImage.open(io.BytesIO(file_bytes))
                meta = (
                    f"Format: {img.format or ext.lstrip('.').upper()}, "
                    f"Size: {img.size[0]}x{img.size[1]}px, "
                    f"Mode: {img.mode}"
                )
                b64 = base64.b64encode(file_bytes).decode("utf-8")
                media_type = _IMAGE_MIME.get(ext, "image/jpeg")
                return {
                    "__vision__": True,
                    "b64": b64,
                    "media_type": media_type,
                    "meta": meta,
                }, None
            except ImportError:
                return None, "Image support requires: pip install Pillow"
            except Exception as exc:
                return None, f"Image processing error: {exc}"

        # SVG
        elif ext == ".svg":
            text = file_bytes.decode("utf-8", errors="replace")
            return f"[SVG Image: {filename}]\n{_truncate(text, 30_000)}", None

        # Archives
        elif ext in (".zip", ".rar"):
            file_storage.seek(0)
            files, err = read_archive(file_storage)
            if err:
                return None, err
            return format_files_for_prompt(files), None

        # Plain text / code
        elif ext in CODE_EXTENSIONS or ext in (".txt", ".md", ".rtf"):
            text = file_bytes.decode("utf-8", errors="replace")
            return _truncate(text, 50_000), None

        # Unknown -- try UTF-8 text
        else:
            try:
                text = file_bytes.decode("utf-8", errors="strict")
                return f"[File: {filename}]\n{_truncate(text, 30_000)}", None
            except UnicodeDecodeError:
                return (
                    None,
                    f"Unsupported file type: {ext}. "
                    "I can analyse text documents, PDFs, Word files, "
                    "spreadsheets, code files, and more.",
                )
    except Exception as exc:
        return None, f"Error reading file: {exc}"

# ===========================================================================
# Prompt formatters
# ===========================================================================

def format_files_for_prompt(files_dict: dict[str, str]) -> str:
    """
    Format a {filename: content} dict into a single prompt-ready string.
    Truncated at ~30 000 characters to stay within model context limits.
    """
    parts = [
        f"--- {path} ---\n{content}"
        for path, content in sorted(files_dict.items())
    ]
    return _truncate("\n\n".join(parts), 30_000)

# ===========================================================================
# Internal helpers
# ===========================================================================

def _truncate(text: str, limit: int) -> str:
    if len(text) > limit:
        return text[:limit] + "\n\n[... truncated due to length ...]"
    return text
