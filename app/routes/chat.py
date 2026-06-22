"""
app/routes/chat.py
==================
Chat API and file-download routes:
  POST /chat — send a message (text or multi-file upload)
  POST /clear — wipe conversation memory for the current mode
  POST /upload-code — upload a zip/rar, get AI-modified code back
  POST /regenerate — regenerate the last assistant response
  GET /download/<filename> — download a generated PPTX
  GET /download-zip/<filename> — download a generated code ZIP
  GET /history — fetch recent conversation history (DB direct)
  GET /memory-sidebar — memory for all modes (DB direct)

v4.0 changes:
  - Multi-image upload support (multiple files in one request)
  - Pollinations.ai PPTX image generation
  - /regenerate endpoint to retry last response
  - Uses ask_ai() for HF-first, Groq-fallback routing
"""

from __future__ import annotations

import base64
import json
import os
import time
from typing import Any

import requests
from flask import (
    Blueprint,
    Response,
    jsonify,
    redirect,
    request,
    send_from_directory,
    session,
    stream_with_context,
    url_for,
)
from pptx import Presentation
from pptx.util import Inches, Pt  # noqa: F401

from app.config.ai_modes import AI_MODES
from app.config.chat_models import DEFAULT_CHAT_MODEL_KEY, apply_chat_model_override
from app.config.settings import Config
from app.services.ai_client import ask_ai, ask_ai_stream, ask_groq_vision
from app.services.external_imports import download_external_file
from app.services.memory import retrieve_relevant_memory
from app.services.search import web_search, format_sources
from app.services.storage import get_user_memory, update_user_memory
from app.utils.artifacts import build_code_download_artifact, extract_requested_artifact_name
from app.utils.files import (
    MAX_UPLOAD_SIZE,
    extract_code_blocks,
    extract_file_content,
    format_files_for_prompt,
    read_archive,
)

chat_bp = Blueprint("chat", __name__)
_cfg = Config()

# ---------------------------------------------------------------------------
# Keywords that indicate the user explicitly wants a downloadable file
# ---------------------------------------------------------------------------
_FILE_GEN_KEYWORDS: list[str] = [
    "create a file", "generate a file", "make a file", "write a file",
    "save as", "download", "export", "create file", "generate file",
    "make file", "write file", "give me the file", "give me a file",
    "create the code", "generate the code", "zip", "package",
    "save the code", "downloadable", "create a project",
    "generate a project", "make a project", "build a project",
]

def _user_wants_file(message: str) -> bool:
    msg_lower = message.lower()
    return any(kw in msg_lower for kw in _FILE_GEN_KEYWORDS)


def _request_model_key() -> str:
    raw = None
    if request.content_type and "multipart/form-data" in (request.content_type or ""):
        raw = request.form.get("model_key")
    else:
        data = request.get_json(silent=True) or {}
        raw = data.get("model_key")

    model_key = (raw or session.get("chat_model_key") or DEFAULT_CHAT_MODEL_KEY).strip()
    session["chat_model_key"] = model_key
    return model_key

# ---------------------------------------------------------------------------
# PPTX generation helpers — v4.0 with Pollinations.ai images
# v2.1: Agentic pipeline — Prompt -> [Image Generator AI] ->
#       [Slide Generator AI] -> (PPTX) Output. All models must come from
#       OpenRouter's FREE tier (Exa AI is web-search only, not used here).
# ---------------------------------------------------------------------------

# Step 2 — "AI Slide Generator": receives the finalized content (with each
# slide flagged for whether an image was generated) and its ONLY job is to
# validate + repackage it into the exact JSON schema for PPTX packing — NO
# editing of titles or bullets. This is a deliberate second OpenRouter
# free-model call, separate from the content-planning model below.
_SLIDE_PACKER_MODE: dict = {
    "name": "AI Slide Generator (internal)",
    "openrouter_model": "meta-llama/llama-3.3-70b-instruct:free",
    "model": "llama-3.3-70b-versatile",
    "system_prompt": (
        "You are the AI Slide Generator, the final packing step of an "
        "agentic presentation pipeline. You receive a JSON presentation "
        "plan whose content has ALREADY been finalized, including a "
        "has_image flag for each slide showing whether an AI-generated "
        "image is attached to it. "
        "Your ONLY job is to validate this content and repackage it into "
        "the exact same JSON schema, ready for PPTX packing. "
        "DO NOT edit, rewrite, shorten, expand, reorder, or rephrase any "
        "titles or bullet points — preserve them EXACTLY as given. "
        "DO NOT add new slides or remove existing slides; the slide count "
        "must stay identical. "
        'Return ONLY a valid JSON object in this exact shape: '
        '{"title": "Presentation Title", "slides": [{"title": "...", '
        '"bullets": ["...", "..."], "has_image": true}]}. '
        "No markdown fences, no explanations, no extra text — JSON only."
    ),
    "temperature": 0.1,
    "max_tokens": 3000,
}


def _generate_pollinations_image(prompt: str, slide_index: int, timestamp: int) -> str | None:
    """
    Download an AI-generated image from Pollinations.ai.
    Returns the local file path or None if it fails.
    """
    try:
        safe_prompt = requests.utils.quote(prompt[:500])
        img_url = f"{_cfg.POLLINATIONS_BASE}/{safe_prompt}?width=1024&height=576&nologo=true&seed={timestamp + slide_index}"
        img_path = f"/tmp/zenith_slide_{timestamp}_{slide_index}.png"

        resp = requests.get(img_url, timeout=30)
        resp.raise_for_status()
        with open(img_path, "wb") as f:
            f.write(resp.content)
        return img_path
    except Exception as exc:
        print(f"[Pollinations error] slide {slide_index}: {exc}")
        return None

def _generate_pptx(ai_response: str, mode: dict | None = None) -> dict | None:
    """
    Parse the AI's JSON response and create a .pptx file in /tmp.
    v4.0: Fetches AI-generated images from Pollinations.ai for each slide.

    Returns a dict with filename, url, and slides keys,
    or None if parsing fails.
    """
    try:
        text = ai_response.strip()
        if text.startswith("```"):
            text = text.split("\n", 1)[1]
            text = text.rsplit("```", 1)[0].strip()
        data = json.loads(text)
    except (json.JSONDecodeError, IndexError):
        return None

    title = data.get("title", "Presentation")
    slides = data.get("slides", [])
    if not slides:
        return None

    prs = Presentation()
    timestamp = int(time.time())
    image_paths: list[str] = []
    slide_titles: list[str] = []

    # ------------------------------------------------------------------
    # v2.1 agentic pipeline, step A — [Image Generator AI]:
    # the content-planning model above already produced an `image_prompt`
    # per slide. Generate the actual AI image for each slide now via
    # Pollinations.ai (free, no key needed) BEFORE handing off to the
    # Slide Generator AI, so it can pack with full knowledge of which
    # slides have images.
    # ------------------------------------------------------------------
    slide_images: list[str | None] = []
    for idx, slide_data in enumerate(slides):
        img_prompt = slide_data.get(
            "image_prompt", f"Professional presentation slide about {slide_data.get('title', 'this topic')}"
        )
        slide_images.append(_generate_pollinations_image(img_prompt, idx, timestamp))

    # ------------------------------------------------------------------
    # v2.1 agentic pipeline, step B — [Slide Generator AI]:
    # a second OpenRouter free-model call whose only job is to validate
    # and repackage the finalized content (titles/bullets unchanged) plus
    # the has_image flags into the exact schema we need. If this step
    # fails for any reason, we fall back to the original step-1 plan so
    # the pipeline still produces a valid PPTX.
    # ------------------------------------------------------------------
    try:
        pack_input = json.dumps({
            "title": title,
            "slides": [
                {
                    "title": s.get("title", "Untitled"),
                    "bullets": s.get("bullets", []),
                    "has_image": slide_images[i] is not None,
                }
                for i, s in enumerate(slides)
            ],
        })
        packed_raw = ask_ai(pack_input, mode=_SLIDE_PACKER_MODE)
        packed_text = packed_raw.strip()
        if packed_text.startswith("```"):
            packed_text = packed_text.split("\n", 1)[1].rsplit("```", 1)[0].strip()
        packed = json.loads(packed_text)
        packed_slides = packed.get("slides", [])
        if packed_slides and len(packed_slides) == len(slides):
            for i, ps in enumerate(packed_slides):
                if ps.get("title"):
                    slides[i]["title"] = ps["title"]
                if ps.get("bullets"):
                    slides[i]["bullets"] = ps["bullets"]
            if packed.get("title"):
                title = packed["title"]
    except Exception as exc:  # pragma: no cover
        print(f"[Slide Generator AI error] {exc} — using original plan")

    # ------------------------------------------------------------------
    # Step C — (PPTX) Output: deterministic packing, no AI involved.
    # ------------------------------------------------------------------
    for idx, slide_data in enumerate(slides):
        slide_title = slide_data.get("title", "Untitled")
        slide_titles.append(slide_title)
        bullets = slide_data.get("bullets", [])

        # Title + content layout
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout

        # Add title
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        tf = title_shape.text_frame
        tf.text = slide_title
        tf.paragraphs[0].font.size = Pt(32)
        tf.paragraphs[0].font.bold = True

        # Add the image generated in step A, if any
        img_path = slide_images[idx]
        if img_path and os.path.exists(img_path):
            try:
                slide.shapes.add_picture(img_path, Inches(5.5), Inches(1.3), width=Inches(4))
                image_paths.append(img_path)
            except Exception:
                img_path = None  # Skip image if it fails

        # Add bullet points
        left = Inches(0.5) if img_path else Inches(0.5)
        width = Inches(4.8) if img_path else Inches(9)
        body_shape = slide.shapes.add_textbox(left, Inches(1.3), width, Inches(5.5))
        btf = body_shape.text_frame
        btf.word_wrap = True

        for j, bullet in enumerate(bullets):
            if j == 0:
                btf.text = bullet
                btf.paragraphs[0].font.size = Pt(18)
                btf.paragraphs[0].space_after = Pt(8)
            else:
                p = btf.add_paragraph()
                p.text = bullet
                p.level = 0
                p.font.size = Pt(18)
                p.space_after = Pt(8)

    filename = f"zenith_ox_{timestamp}.pptx"
    filepath = f"/tmp/{filename}"
    prs.save(filepath)

    return {
        "filename": filename,
        "path": filepath,
        "url": f"/download/{filename}",
        "slides": slide_titles,
    }

# ---------------------------------------------------------------------------
# Multi-file upload processor — v4.0
# ---------------------------------------------------------------------------

def _process_uploads(files: list) -> tuple[str, list[dict]]:
    """
    Process multiple uploaded files.
    Returns (combined_context, vision_payloads).
    vision_payloads are dicts with __vision__ for image files.
    """
    contexts: list[str] = []
    visions: list[dict] = []

    for f in files:
        if not f or not f.filename:
            continue
        content, error = extract_file_content(f)
        if error:
            continue
        if isinstance(content, dict) and content.get("__vision__"):
            visions.append(content)
        elif content:
            contexts.append(
                f"\n\n--- Uploaded File: {f.filename} ---\n"
                f"{content}\n--- End of File ---\n"
            )
    return "\n".join(contexts), visions

# ===========================================================================
# Routes
# ===========================================================================

@chat_bp.route("/chat", methods=["POST"])
def chat():
    if "user_id" not in session:
        return jsonify({"ok": False, "error": "Not authenticated."}), 401

    user_id = session["user_id"]
    mode_key = session.get("ai_mode", "researcher")
    mode = apply_chat_model_override(
        AI_MODES.get(mode_key, AI_MODES["researcher"]),
        _request_model_key(),
    )

    # Parse request (multipart for uploads, JSON otherwise)
    uploaded_files: list = []
    message = ""

    if request.content_type and "multipart/form-data" in request.content_type:
        message = (request.form.get("message") or "").strip()
        # v4.0: support multiple files
        uploaded_files = request.files.getlist("files") or []
        if not uploaded_files:
            single = request.files.get("file")
            if single:
                uploaded_files = [single]
    else:
        data = request.get_json(silent=True) or {}
        message = (data.get("message") or "").strip()

    if not message and not uploaded_files:
        return jsonify({"ok": False, "error": "Empty message."}), 400

    # Load memory
    memory_key = f"{user_id}:{mode_key}"
    recent_history = get_user_memory(memory_key)

    file_context = ""
    vision_payloads: list[dict] = []

    # v4.0: Process multiple uploads
    if uploaded_files:
        file_context, vision_payloads = _process_uploads(uploaded_files)

    # Handle vision images first (Groq Vision for each image)
    vision_answers: list[str] = []
    if vision_payloads:
        for vp in vision_payloads:
            try:
                ans = ask_groq_vision(
                    message=message or "Please analyse this image and describe what you see.",
                    b64_image=vp["b64"],
                    media_type=vp["media_type"],
                    meta=vp.get("meta", ""),
                    mode=mode,
                    recent_history=recent_history,
                )
                vision_answers.append(ans)
            except Exception as exc:
                vision_answers.append(f"[Image analysis error: {exc}]")

        # If we only had images (no text files), combine vision answers
        if not file_context and vision_answers:
            vision_summary = "\n\n".join(
                f"[Image {i+1} Analysis]: {ans}" for i, ans in enumerate(vision_answers)
            )
            full_message = (
                f"{message}\n\n{vision_summary}" if message
                else vision_summary
            )
            # Save to memory and return
            try:
                update_user_memory(
                    memory_key,
                    "user",
                    message or f"[Uploaded {len(vision_payloads)} image(s)]",
                )
                update_user_memory(memory_key, "assistant", vision_summary)
            except Exception:
                pass
            return jsonify({"ok": True, "response": vision_summary})

    # Combine everything into the full message
    parts: list[str] = []
    if message:
        parts.append(message)
    if file_context:
        parts.append(file_context)
    if vision_answers:
        parts.append("\n\n".join(
            f"[Image {i+1}]: {ans}" for i, ans in enumerate(vision_answers)
        ))

    full_message = "\n\n".join(parts).strip()
    if not full_message:
        return jsonify({"ok": False, "error": "Empty message."}), 400

    vector_mem = retrieve_relevant_memory(user_id, mode_key, message or "file analysis")
    search_result = web_search(message) if mode.get("uses_web_search") and message else {}
    web_ctx = search_result.get("context", "")
    sources = search_result.get("sources", [])

    try:
        # v4.0: Use ask_ai() for HF-first, Groq-fallback routing
        answer = ask_ai(full_message, vector_mem, web_ctx, mode, recent_history=recent_history)
    except Exception as exc:
        return jsonify({"ok": False, "error": f"AI error: {exc}"}), 500

    # v2.1: After a deep-research web search, sources MUST be stated.
    if web_ctx and sources:
        answer += format_sources(sources)

    # PPTX special handler — v4.0 with images
    if mode.get("special_handler") == "pptx":
        result = _generate_pptx(answer, mode)
        if result:
            slide_list = "".join(
                f"  {i}. {t}\n" for i, t in enumerate(result["slides"], 1)
            )
            summary = (
                f"Your presentation is ready!\n\nSlides:\n{slide_list}\n"
                "Click the download button below to save your file."
            )
            update_user_memory(memory_key, "user", message)
            update_user_memory(memory_key, "assistant", summary)
            return jsonify({
                "ok": True,
                "response": summary,
                "download_url": result["url"],
                "download_name": result["filename"],
            })

    # Developer mode: ZIP on explicit file request
    if mode_key == "developer" and _user_wants_file(message):
        code_blocks = extract_code_blocks(answer)
        if code_blocks:
            zip_info = build_code_download_artifact(
                code_blocks,
                requested_name=extract_requested_artifact_name(message),
            )
            if zip_info:
                update_user_memory(memory_key, "user", message)
                update_user_memory(memory_key, "assistant", answer)
                return jsonify({
                    "ok": True,
                    "response": answer,
                    "download_url": zip_info["url"],
                    "download_name": zip_info["filename"],
                })

    # Standard response
    user_content = message or f"[Uploaded: {', '.join(f.filename for f in uploaded_files if f)}]"
    update_user_memory(memory_key, "user", user_content)
    update_user_memory(memory_key, "assistant", answer)
    return jsonify({"ok": True, "response": answer})


@chat_bp.route("/chat/stream", methods=["POST"])
def chat_stream():
    """
    v2.1 — SSE streaming chat ("AI text tracker").

    POST JSON {"message": "..."} (text-only — file uploads still go
    through /chat). Streams the AI's response as Server-Sent Events:

      data: {"delta": "..."}        — one or more text chunks
      data: {"error": "..."}        — something went wrong mid-stream
      data: {"done": true, ...}      — final event; may include
                                        download_url / download_name for
                                        Developer-mode code ZIPs.

    Slides Generator (pptx) is intentionally NOT streamed here — its raw
    JSON output isn't meant to be shown to the user, so it keeps using the
    non-streaming /chat endpoint.
    """
    if "user_id" not in session:
        return jsonify({"ok": False, "error": "Not authenticated."}), 401

    user_id = session["user_id"]
    mode_key = session.get("ai_mode", "researcher")
    mode = apply_chat_model_override(
        AI_MODES.get(mode_key, AI_MODES["researcher"]),
        _request_model_key(),
    )

    if mode.get("special_handler") == "pptx":
        return jsonify({"ok": False, "error": "Slides Generator does not support streaming."}), 400

    data = request.get_json(silent=True) or {}
    message = (data.get("message") or "").strip()
    if not message:
        return jsonify({"ok": False, "error": "Empty message."}), 400

    memory_key = f"{user_id}:{mode_key}"
    recent_history = get_user_memory(memory_key)

    vector_mem = retrieve_relevant_memory(user_id, mode_key, message)
    search_result = web_search(message) if mode.get("uses_web_search") else {}
    web_ctx = search_result.get("context", "")
    sources = search_result.get("sources", [])

    def generate():
        chunks: list[str] = []
        try:
            for piece in ask_ai_stream(message, vector_mem, web_ctx, mode, recent_history=recent_history):
                chunks.append(piece)
                yield f"data: {json.dumps({'delta': piece})}\n\n"
        except Exception as exc:
            yield f"data: {json.dumps({'error': f'AI error: {exc}'})}\n\n"
            return

        answer = "".join(chunks)

        # v2.1: After a deep-research web search, sources MUST be stated.
        if web_ctx and sources:
            src_md = format_sources(sources)
            answer += src_md
            yield f"data: {json.dumps({'delta': src_md})}\n\n"

        done_payload: dict[str, Any] = {"done": True}

        # Developer mode: ZIP on explicit file request
        if mode_key == "developer" and _user_wants_file(message):
            code_blocks = extract_code_blocks(answer)
            if code_blocks:
                zip_info = build_code_download_artifact(
                    code_blocks,
                    requested_name=extract_requested_artifact_name(message),
                )
                if zip_info:
                    done_payload["download_url"] = zip_info["url"]
                    done_payload["download_name"] = zip_info["filename"]

        try:
            update_user_memory(memory_key, "user", message)
            update_user_memory(memory_key, "assistant", answer)
        except Exception:
            pass

        yield f"data: {json.dumps(done_payload)}\n\n"

    return Response(
        stream_with_context(generate()),
        mimetype="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "X-Accel-Buffering": "no",  # disable proxy buffering (Vercel/Nginx)
        },
    )

@chat_bp.route("/regenerate", methods=["POST"])
def regenerate():
    """
    v4.0: Regenerate the last assistant response.
    Re-sends the last user message to the AI and returns a new response.
    """
    if "user_id" not in session:
        return jsonify({"ok": False, "error": "Not authenticated."}), 401

    user_id = session["user_id"]
    mode_key = session.get("ai_mode", "researcher")
    mode = apply_chat_model_override(
        AI_MODES.get(mode_key, AI_MODES["researcher"]),
        _request_model_key(),
    )

    # Get the last user message from memory
    memory_key = f"{user_id}:{mode_key}"
    recent_history = get_user_memory(memory_key)

    if not recent_history:
        return jsonify({"ok": False, "error": "No conversation history to regenerate from."}), 400

    # Find the last exchange (user message followed by assistant response)
    last_user_msg = ""
    last_assistant_idx = -1
    for i in range(len(recent_history) - 1, -1, -1):
        if recent_history[i]["role"] == "assistant" and last_assistant_idx == -1:
            last_assistant_idx = i
        elif recent_history[i]["role"] == "user" and last_assistant_idx != -1:
            last_user_msg = recent_history[i]["content"]
            break

    if not last_user_msg:
        # Try just the last user message
        for i in range(len(recent_history) - 1, -1, -1):
            if recent_history[i]["role"] == "user":
                last_user_msg = recent_history[i]["content"]
                break

    if not last_user_msg:
        return jsonify({"ok": False, "error": "No user message found to regenerate from."}), 400

    # Strip file context markers for regeneration
    clean_msg = last_user_msg
    if "--- Uploaded File:" in clean_msg:
        clean_msg = clean_msg.split("--- Uploaded File:")[0].strip()

    vector_mem = retrieve_relevant_memory(user_id, mode_key, clean_msg)
    search_result = web_search(clean_msg) if mode.get("uses_web_search") else {}
    web_ctx = search_result.get("context", "")
    sources = search_result.get("sources", [])

    # Build history without the last assistant response
    trimmed_history = recent_history[:last_assistant_idx] if last_assistant_idx > 0 else recent_history[:-1]

    try:
        # Use a slightly higher temperature for variation
        regen_mode = dict(mode)
        regen_mode["temperature"] = min(regen_mode.get("temperature", 0.7) + 0.1, 1.0)
        answer = ask_ai(clean_msg, vector_mem, web_ctx, regen_mode, recent_history=trimmed_history)
    except Exception as exc:
        return jsonify({"ok": False, "error": f"Regeneration error: {exc}"}), 500

    # v2.1: After a deep-research web search, sources MUST be stated.
    if web_ctx and sources:
        answer += format_sources(sources)

    # Update memory with the new response (replace last assistant entry)
    update_user_memory(memory_key, "assistant", answer)

    # PPTX special handler on regenerate
    if mode.get("special_handler") == "pptx":
        result = _generate_pptx(answer, mode)
        if result:
            slide_list = "".join(
                f"  {i}. {t}\n" for i, t in enumerate(result["slides"], 1)
            )
            summary = (
                f"Your presentation is ready!\n\nSlides:\n{slide_list}\n"
                "Click the download button below to save your file."
            )
            return jsonify({
                "ok": True,
                "response": summary,
                "download_url": result["url"],
                "download_name": result["filename"],
            })

    # Developer mode: ZIP on explicit file request
    if mode_key == "developer" and _user_wants_file(clean_msg):
        code_blocks = extract_code_blocks(answer)
        if code_blocks:
            zip_info = build_code_download_artifact(
                code_blocks,
                requested_name=extract_requested_artifact_name(clean_msg),
            )
            if zip_info:
                return jsonify({
                    "ok": True,
                    "response": answer,
                    "download_url": zip_info["url"],
                    "download_name": zip_info["filename"],
                })

    return jsonify({"ok": True, "response": answer})

@chat_bp.route("/clear", methods=["POST"])
def clear():
    if "user_id" not in session:
        return jsonify({"ok": False, "error": "Not authenticated."}), 401

    user_id = session["user_id"]
    mode_key = session.get("ai_mode", "researcher")

    try:
        from app.services.db import db_clear_conversation
        db_clear_conversation(user_id, mode_key)
    except RuntimeError as exc:
        return jsonify({"ok": False, "error": f"Could not clear memory: {exc}"}), 500

    return jsonify({"ok": True, "message": "Memory cleared for this mode."})

@chat_bp.route("/upload-code", methods=["POST"])
def upload_code():
    """Upload a ZIP/RAR project; AI analyses or modifies it and returns a new ZIP."""
    if "user_id" not in session:
        return jsonify({"ok": False, "error": "Not authenticated."}), 401

    # Code-project ZIP/RAR upload is a Developer-mode-only feature.
    if session.get("ai_mode") != "developer":
        return jsonify({
            "ok": False,
            "error": "Code project upload is only available in Developer mode. "
                     "Switch to Developer mode and try again.",
        }), 403

    if "file" not in request.files:
        return jsonify({"ok": False, "error": "No file uploaded."}), 400

    uploaded = request.files["file"]
    if not uploaded.filename:
        return jsonify({"ok": False, "error": "Empty filename."}), 400

    instruction = (
        request.form.get("message", "").strip()
        or "Analyse this code and suggest improvements."
    )

    files_dict, error = read_archive(uploaded)
    if error:
        return jsonify({"ok": False, "error": error}), 400

    files_text = format_files_for_prompt(files_dict)
    file_list = ", ".join(files_dict.keys())

    user_id = session["user_id"]
    mode_key = "developer"
    mode = apply_chat_model_override(
        AI_MODES["developer"],
        _request_model_key(),
    )
    memory_key = f"{user_id}:{mode_key}"

    upload_prompt = (
        f"The user uploaded a code project with these files: {file_list}\n\n"
        f"Here are the file contents:\n\n{files_text}\n\n"
        f"User instruction: {instruction}\n\n"
        "IMPORTANT: When returning modified code, wrap each file in markdown "
        "code fences with the language name. If modifying multiple files, "
        "include ALL of them."
    )

    recent_history = get_user_memory(memory_key)
    vector_mem = retrieve_relevant_memory(user_id, mode_key, instruction)

    try:
        answer = ask_ai(upload_prompt, vector_mem, "", mode, recent_history=recent_history)
    except Exception as exc:
        return jsonify({"ok": False, "error": f"AI error: {exc}"}), 500

    code_blocks = extract_code_blocks(answer)
    response_data = {"ok": True, "response": answer}

    if code_blocks:
        zip_info = build_code_download_artifact(
            code_blocks,
            requested_name=extract_requested_artifact_name(instruction),
        )
        if zip_info:
            response_data["download_url"] = zip_info["url"]
            response_data["download_name"] = zip_info["filename"]

    update_user_memory(memory_key, "user", f"[Uploaded: {file_list}] {instruction}")
    update_user_memory(memory_key, "assistant", answer)
    return jsonify(response_data)



@chat_bp.route("/import-external", methods=["POST"])
def import_external():
    if "user_id" not in session:
        return jsonify({"ok": False, "error": "Not authenticated."}), 401

    data = request.get_json(silent=True) or {}
    provider = (data.get("provider") or "").strip().lower()
    url = (data.get("url") or "").strip()
    if not provider or not url:
        return jsonify({"ok": False, "error": "Provider and URL are required."}), 400

    try:
        payload, filename, content_type = download_external_file(provider, url, MAX_UPLOAD_SIZE)
    except Exception as exc:
        return jsonify({"ok": False, "error": str(exc)}), 400

    return jsonify({
        "ok": True,
        "provider": provider,
        "filename": filename,
        "content_type": content_type,
        "size_bytes": len(payload),
        "data_b64": base64.b64encode(payload).decode("utf-8"),
    })

@chat_bp.route("/download-generated/<path:filename>")
def download_generated(filename: str):
    if "user_id" not in session:
        return redirect(url_for("auth.login_page"))
    safe_name = os.path.basename(filename)
    filepath = os.path.join("/tmp", safe_name)
    if not os.path.exists(filepath):
        return "File not found or expired", 404
    return send_from_directory("/tmp", safe_name, as_attachment=True, download_name=safe_name)

@chat_bp.route("/download/<filename>")
def download_file(filename: str):
    if "user_id" not in session:
        return redirect(url_for("auth.login_page"))
    if not filename.endswith(".pptx"):
        return "Invalid file type", 400
    filepath = os.path.join("/tmp", filename)
    if not os.path.exists(filepath):
        return "File not found or expired", 404
    return send_from_directory("/tmp", filename, as_attachment=True)

@chat_bp.route("/download-zip/<filename>")
def download_zip(filename: str):
    if "user_id" not in session:
        return redirect(url_for("auth.login_page"))
    if not filename.endswith(".zip"):
        return "Invalid file type", 400
    filepath = os.path.join("/tmp", filename)
    if not os.path.exists(filepath):
        return "File not found or expired", 404
    return send_from_directory("/tmp", filename, as_attachment=True)

@chat_bp.route("/history")
def history():
    if "user_id" not in session:
        return jsonify({"ok": False, "error": "Not authenticated."}), 401

    user_id = session["user_id"]
    mode_key = session.get("ai_mode", "researcher")

    try:
        from app.services.db import db_get_conversation
        messages = db_get_conversation(user_id, mode_key, limit=30)
    except RuntimeError as exc:
        return jsonify({"ok": False, "error": f"Could not load history: {exc}"}), 500

    return jsonify({"ok": True, "messages": messages})

@chat_bp.route("/memory-sidebar")
def memory_sidebar():
    """
    Return structured memory for the sidebar — grouped by mode, all modes.
    Works when logged out too (returns empty per-mode dicts).
    """
    user_id = session.get("user_id")
    if not user_id:
        modes_out = {key: [] for key in AI_MODES}
        return jsonify({"ok": True, "modes": modes_out, "logged_in": False})

    try:
        from app.services.db import db_get_all_modes_memory
        all_modes = db_get_all_modes_memory(user_id)
    except RuntimeError as exc:
        return jsonify({"ok": False, "error": f"Could not load memory: {exc}"}), 500

    modes_out = {}
    for mode_key in AI_MODES:
        raw = all_modes.get(mode_key, [])
        exchanges = []
        for i, msg in enumerate(raw):
            if msg["role"] == "user":
                reply = ""
                if i + 1 < len(raw) and raw[i + 1]["role"] == "assistant":
                    reply = raw[i + 1]["content"][:160]
                exchanges.append({
                    "user": msg["content"][:120],
                    "assistant": reply,
                })
        modes_out[mode_key] = exchanges[-20:]  # last 20 exchanges per mode

    return jsonify({"ok": True, "modes": modes_out, "logged_in": True})

@chat_bp.route("/export-chat", methods=["POST"])
def export_chat():
    """
    v4.0: Export the current chat as various formats.
    Supports: json, txt, md
    """
    if "user_id" not in session:
        return jsonify({"ok": False, "error": "Not authenticated."}), 401

    data = request.get_json(silent=True) or {}
    fmt = data.get("format", "md")
    messages = data.get("messages", [])
    mode = data.get("mode", "researcher")

    if not messages:
        return jsonify({"ok": False, "error": "No messages to export."}), 400

    timestamp = int(time.time())

    if fmt == "json":
        export_data = {
            "title": f"Zenith OX Chat — {mode}",
            "mode": mode,
            "exported_at": timestamp,
            "messages": messages,
        }
        filename = f"zenith_chat_{timestamp}.json"
        filepath = f"/tmp/{filename}"
        with open(filepath, "w") as f:
            json.dump(export_data, f, indent=2)
        return jsonify({"ok": True, "url": f"/download-export/{filename}"})

    elif fmt in ("txt", "md"):
        ext = fmt
        lines = [f"# Zenith OX Chat — {mode}", f"Exported: {timestamp}", "=" * 50, ""]
        for msg in messages:
            role = msg.get("role", "unknown").upper()
            content = msg.get("content", "")
            lines.append(f"**{role}:**\n{content}\n")
        content = "\n".join(lines)
        filename = f"zenith_chat_{timestamp}.{ext}"
        filepath = f"/tmp/{filename}"
        with open(filepath, "w") as f:
            f.write(content)
        return jsonify({"ok": True, "url": f"/download-export/{filename}"})

    return jsonify({"ok": False, "error": f"Unsupported format: {fmt}"}), 400

@chat_bp.route("/download-export/<filename>")
def download_export(filename: str):
    if "user_id" not in session:
        return redirect(url_for("auth.login_page"))
    filepath = os.path.join("/tmp", filename)
    if not os.path.exists(filepath):
        return "File not found or expired", 404
    return send_from_directory("/tmp", filename, as_attachment=True)
